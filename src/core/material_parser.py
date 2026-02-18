import fitz  # PyMuPDF
from pptx import Presentation
import os

class MaterialParser:
    def __init__(self):
        pass

    def parse_pdf(self, pdf_path):
        """
        Extracts text from a PDF file.
        """
        print(f"Parsing PDF: {pdf_path}")
        text = ""
        with fitz.open(pdf_path) as doc:
            for page in doc:
                text += page.get_text() + "\n"
        return text

    def parse_pptx(self, pptx_path):
        """
        Extracts text from a PPTX file.
        """
        print(f"Parsing PPTX: {pptx_path}")
        text = ""
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def parse_material(self, file_path):
        """
        Detects file type and parses accordingly.
        """
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            return self.parse_pdf(file_path)
        elif ext == ".pptx":
            return self.parse_pptx(file_path)
        else:
            print(f"Unsupported material format: {ext}")
            return None

if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        parser = MaterialParser()
        content = parser.parse_material(sys.argv[1])
        if content:
            print("--- Extracted Content (first 500 chars) ---")
            print(content[:500])
    else:
        print("Usage: python material_parser.py <path_to_material>")
